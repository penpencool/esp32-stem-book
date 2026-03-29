#!/usr/bin/env python3
"""Demo: Create 3 slides about ESP32 using python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation (16:9 vertical for TikTok)
prs = Presentation()
prs.slide_width = Inches(9)
prs.slide_height = Inches(16)

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(255, 136, 0)

def add_title_slide(prs, title, subtitle):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(9), Inches(16)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    """Slide 2: Content slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(9), Inches(16)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shape.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(9), Inches(2.5)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(7.6), Inches(12))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(18)

def add_code_slide(prs, title, code):
    """Slide 3: Code example slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(9), Inches(16)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # Code block
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(12))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(18)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(0, 255, 136)

# Create 3 slides
# Slide 1: Title
add_title_slide(
    prs,
    "ESP32-C3",
    "บอร์ดไıkสำหรับ IoT และ Embedded Systems"
)

# Slide 2: What is ESP32
add_content_slide(
    prs,
    "ESP32 คืออะไร?",
    [
        "ไมโครคอนโทรลเลอร์ที่ทรงพลัง",
        "มี WiFi และ Bluetooth ในตัว",
        "ราคาถูก (ประมาณ 100-200 บาท)",
        "ใช้งานง่าย มี Arduino IDE รองรับ",
        "เหมาะสำหรับโปรเจกต์ IoT"
    ]
)

# Slide 3: Code example
add_code_slide(
    prs,
    "ตัวอย่าง: Hello World",
    """void setup() {
  Serial.begin(115200);
  Serial.println("Hello ESP32!");
}

void loop() {
  // ไม่ต้องทำอะไร
}"""
)

# Save
output_path = "/home/maxtic/esp32-stem-book/demo_slides.pptx"
prs.save(output_path)
print(f"✅ Created: {output_path}")
