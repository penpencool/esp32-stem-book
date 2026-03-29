#!/usr/bin/env python3
"""
สคริปต์สร้าง Slide บทที่ 1: AI คืออะไร?
Design: Dark theme, ESP32 style
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(13, 27, 42)       # #0D1B2A
CYAN = RGBColor(0, 217, 255)           # #00D9FF
PINK = RGBColor(255, 107, 138)        # #FF6B8A
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(180, 180, 200)
GOLD = RGBColor(255, 215, 0)
GREEN = RGBColor(0, 255, 136)

def add_slide_background(slide, prs):
    """เพิ่มพื้นหลังสีเข้ม"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

def add_title_bar(slide, text, prs):
    """เพิ่ม title bar สี gradient"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BG

def add_emoji_title(slide, emoji, text, top):
    """เพิ่มหัวข้อพร้อม emoji"""
    title = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {text}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CYAN

def add_bullet_list(slide, items, top, font_size=22):
    """เพิ่ม bullet list"""
    content = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), Inches(4))
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

def add_table(slide, headers, rows, top):
    """เพิ่มตาราง"""
    cols = len(headers)
    table = slide.shapes.add_table(len(rows)+1, cols, Inches(0.5), top, Inches(12), Inches(3)).table
    
    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
    
    # Rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE

def add_code_box(slide, code, top):
    """เพิ่มกล่องโค้ด"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top,
        Inches(12), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(25, 25, 35)
    box.line.color.rgb = GREEN
    box.line.width = Pt(2)
    
    code_text = slide.shapes.add_textbox(Inches(0.8), top + Inches(0.2), Inches(11.5), Inches(2))
    tf = code_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(16)
    p.font.name = "Courier New"
    p.font.color.rgb = GREEN

def add_horizontal_line(slide, top, color=CYAN):
    """เพิ่มเส้นแนวนอน"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), top,
        Inches(12), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_cta_box(slide, text, top):
    """เพิ่มกล่อง CTA"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), top,
        Inches(9), Inches(1)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PINK
    box.line.fill.background()
    
    cta = slide.shapes.add_textbox(Inches(2.5), top + Inches(0.25), Inches(8), Inches(0.6))
    tf = cta.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def create_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: HOOK ==========
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s1, prs)
    
    # Hook text
    hook_box = s1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(2))
    tf = hook_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Stop scrolling — ถ้าคุณใช้ AI ผิดวิธี"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER
    
    # Sub text
    sub = s1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "มันอาจไม่ใช่ความผิดของ AI\nแต่เป็นเพราะเรา 'ถามผิดวิธี' ต่างหาก!"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Topic
    topic = s1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    tf = topic.text_frame
    p = tf.paragraphs[0]
    p.text = "บทที่ 1: AI คืออะไร? ทำไมต้องรู้!"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 2: AI อยู่รอบตัว ==========
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s2, prs)
    add_title_bar(s2, "🤖 AI อยู่รอบตัวเราแล้ว!", prs)
    
    items = [
        "🌤️ ถาม Google ว่าอากาศเป็นยังไง → AI ประมวลผลคำถาม",
        "📺 ดู YouTube แล้วเห็นคำแนะนำวิดีโอ → AI วิเคราะห์ความชอบ",
        "📸 ถ่ายรูปแล้วมือถือตรวจจับใบหน้า → AI จดจำใบหน้า",
        "🗣️ สั่งงาน Siri หรือ Google Assistant → AI เข้าใจคำสั่งเสียง"
    ]
    add_bullet_list(s2, items, Inches(1.8), font_size=24)
    
    # ========== SLIDE 3: AI คืออะไร? ==========
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s3, prs)
    add_title_bar(s3, "🤖 AI คืออะไร?", prs)
    
    add_emoji_title(s3, "📚", "Artificial Intelligence — ปัญญาประดิษฐ์", Inches(1.5))
    
    desc = s3.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12), Inches(1))
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "โปรแกรมคอมพิวเตอร์ที่ถูกสอนให้ 'คิด' และ 'เข้าใจ' ได้เหมือนมนุษย์ แต่เร็วกว่าเยอะ!"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    
    add_horizontal_line(s3, Inches(3.3))
    
    headers = ["ความสามารถ", "ความหมาย"]
    rows = [
        ["เรียนรู้", "จากข้อมูล (Learning)"],
        ["เข้าใจ", "ภาษามนุษย์ (Understanding)"],
        ["ตัดสินใจ", "จากข้อมูลที่ได้รับ (Decision Making)"],
        ["สร้าง", "ผลลัพธ์ใหม่ๆ (Generating)"]
    ]
    add_table(s3, headers, rows, Inches(3.6))
    
    # ========== SLIDE 4: AI ทำงานยังไง? ==========
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s4, prs)
    add_title_bar(s4, "🧠 AI ทำงานยังไง?", prs)
    
    steps = [
        "1️⃣ รับคำถามของเรา",
        "2️⃣ ค้นหาข้อมูลที่เกี่ยวข้อง",
        "3️⃣ ประมวลผล",
        "4️⃣ ตอบกลับ"
    ]
    add_bullet_list(s4, steps, Inches(1.8), font_size=28)
    
    add_horizontal_line(s4, Inches(4.5), PINK)
    
    note = s4.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(1.5))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ สิ่งสำคัญ:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    p2 = tf.add_paragraph()
    p2.text = "AI ไม่ได้ 'คิด' เหมือนมนุษย์ — มันเป็นโปรแกรมที่ทำนายคำตอบที่ 'น่าจะถูกต้อง'"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    
    # ========== SLIDE 5: Generative AI ==========
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s5, prs)
    add_title_bar(s5, "✨ Generative AI — AI ที่สร้างงานได้", prs)
    
    sub5 = s5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = sub5.text_frame
    p = tf.paragraphs[0]
    p.text = "AI ที่เราจะใช้กันในหนังสือเล่มนี้ สามารถ..."
    p.font.size = Pt(22)
    p.font.color.rgb = GRAY
    
    items5 = [
        "✍️ เขียนข้อความ, บทความ, เรื่องสั้น",
        "💻 เขียนโค้ดโปรแกรม",
        "🎨 วาดรูปจากคำอธิบาย",
        "🔍 อธิบายเรื่องยากให้เข้าใจง่าย"
    ]
    add_bullet_list(s5, items5, Inches(2.3), font_size=24)
    
    # ========== SLIDE 6: AI ช่วยเขียนโค้ด ==========
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s6, prs)
    add_title_bar(s6, "⚡ AI ช่วยเขียนโค้ดยังไง?", prs)
    
    headers6 = ["งาน", "ถาม AI ว่า..."]
    rows6 = [
        ["📖 อธิบายโค้ด", "\"โค้ดนี้ทำงานยังไง?\""],
        ["✍️ เขียนโค้ดใหม่", "\"เขียนโค้ดให้ LED กระพริบได้มั้ย?\""],
        ["🔧 แก้ไขโค้ด", "\"โค้ดนี้ผิดตรงไหน?\""],
        ["📝 เพิ่มเติมโค้ด", "\"เพิ่มฟังก์ชันส่ง Wi-Fi ได้มั้ย?\""]
    ]
    add_table(s6, headers6, rows6, Inches(1.5))
    
    note6 = s6.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1))
    tf = note6.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 AI จึงเป็น 'เพื่อนฉลาด' ที่ช่วยเราสร้างโปรเจกต์ได้เร็วขึ้นมาก!"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # ========== SLIDE 7: เครื่องมือที่ใช้ ==========
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s7, prs)
    add_title_bar(s7, "🛠️ เครื่องมือ AI ที่จะใช้", prs)
    
    headers7 = ["เครื่องมือ", "ใช้ทำอะไร"]
    rows7 = [
        ["ChatGPT", "ถามตอบทุกเรื่อง, เขียนโค้ด"],
        ["GitHub Copilot", "เขียนโค้ดแบบ auto-complete"]
    ]
    add_table(s7, headers7, rows7, Inches(2))
    
    note7 = s7.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(1))
    tf = note7.text_frame
    p = tf.paragraphs[0]
    p.text = "📝 ทั้งสองฟรี! (มี free tier)"
    p.font.size = Pt(20)
    p.font.color.rgb = GOLD
    
    # ========== SLIDE 8: CTA ==========
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s8, prs)
    
    cta_box = s8.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💬 ถามเพื่อนๆ หน่อย"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "พี่เคยใช้ AI ช่วยเขียนโค้ดมั้ย?\nบอกในคอมเมนต์มาเลย! 👇"
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    add_cta_box(s8, "Follow ก่อนนะ เดี๋ยวบทต่อไปสนุกกว่านี้!", Inches(5))
    
    # ========== SLIDE 9: สรุป ==========
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(s9, prs)
    add_title_bar(s9, "📋 สรุปบทที่ 1", prs)
    
    summary = [
        "🤖 AI คืออะไร → โปรแกรมที่เรียนรู้และสร้างงานได้",
        "✨ Generative AI → AI ที่เขียนข้อความ, โค้ด, รูปได้",
        "⚡ AI ช่วยเขียนโค้ด → อธิบาย, เขียน, แก้ไขโค้ด",
        "🛠️ เครื่องมือที่จะใช้ → ChatGPT + GitHub Copilot"
    ]
    add_bullet_list(s9, summary, Inches(1.8), font_size=24)
    
    # Save
    output = "/home/maxtic/esp32-stem-book/chapter01/slides_ch01.pptx"
    prs.save(output)
    print(f"✅ Created: {output}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output

if __name__ == "__main__":
    create_slides()
