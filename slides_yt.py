#!/usr/bin/env python3
"""Create beautiful YouTube slides with python-pptx
Landscape 16:9 format with modern design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Modern color palette
COLORS = {
    'dark_bg': RGBColor(15, 15, 25),      # Dark navy
    'gradient_start': RGBColor(30, 30, 60),
    'gradient_end': RGBColor(60, 20, 90),
    'primary': RGBColor(0, 212, 255),      # Cyan
    'secondary': RGBColor(255, 107, 138),  # Pink
    'accent': RGBColor(255, 215, 0),       # Gold
    'text_light': RGBColor(255, 255, 255),
    'text_gray': RGBColor(180, 180, 200),
    'code_bg': RGBColor(25, 25, 35),
    'green': RGBColor(0, 255, 136),
}

def set_gradient(shape, color1, color2):
    """Set gradient fill for a shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_decoration_circle(slide, left, top, size, color, alpha=0.3):
    """Add decorative circle element"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def create_youtube_slide_1(prs):
    """Slide 1: Eye-catching Hook/Intro"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Full dark gradient background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_gradient(bg, COLORS['gradient_start'], COLORS['gradient_end'])
    bg.line.fill.background()
    
    # Decorative circles
    add_decoration_circle(slide, Inches(9), Inches(-1), Inches(4), COLORS['primary'], 0.15)
    add_decoration_circle(slide, Inches(10), Inches(5), Inches(3), COLORS['secondary'], 0.1)
    add_decoration_circle(slide, Inches(-1), Inches(5), Inches(2), COLORS['accent'], 0.1)
    
    # Glowing line accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(3.2),
        Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🤖 ESP32-C3"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(10), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "บอร์ด IoT ยอดนิยม ราคาถูก ทำงานได้ทุกอย่าง!"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['text_gray']
    
    # Bottom tagline
    tag_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(10), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "▶ เริ่มเรียนรู้วันนี้"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['primary']
    p.font.bold = True

def create_youtube_slide_2(prs):
    """Slide 2: Content - What is ESP32"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark_bg']
    bg.line.fill.background()
    
    # Header bar with gradient
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.8)
    )
    set_gradient(header, COLORS['gradient_start'], COLORS['gradient_end'])
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 ESP32 คืออะไร?"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    
    # Content area - Left side (text)
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(6), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("💡", "ไมโครคอนโทรลเลอร์ทรงพลัง"),
        ("📶", "มี WiFi + Bluetooth ในตัว"),
        ("💰", "ราคาถูกมาก (100-200 บาท)"),
        ("🔧", "ใช้งานง่าย มี Arduino IDE"),
        ("🌐", "เหมาะกับงาน IoT ทุกประเภท"),
    ]
    
    for i, (emoji, text) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{emoji}  {text}"
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['text_light']
        p.space_after = Pt(16)
    
    # Right side - decorative ESP32 image placeholder
    img_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(2.2),
        Inches(4.5), Inches(4)
    )
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
    img_box.line.color.rgb = COLORS['primary']
    img_box.line.width = Pt(2)
    
    # Placeholder text
    placeholder = slide.shapes.add_textbox(Inches(8.3), Inches(3.8), Inches(4), Inches(1))
    tf = placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "📷 รูปบอร์ด ESP32"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_gray']
    p.alignment = PP_ALIGN.CENTER

def create_youtube_slide_3(prs):
    """Slide 3: Code Example"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['code_bg']
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💻 โค้ดตัวอย่าง: Hello World"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    # Code block background
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(8), Inches(5.5)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(20, 20, 30)
    code_bg.line.color.rgb = COLORS['green']
    code_bg.line.width = Pt(2)
    
    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.5), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    code = """void setup() {
  Serial.begin(115200);
  Serial.println("Hello ESP32!");
}

void loop() {
  // พิมพ์ทุก 1 วินาที
  Serial.println("Running...");
  delay(1000);
}"""
    
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(20)
    p.font.name = "Courier New"
    p.font.color.rgb = COLORS['green']
    
    # Side info panel
    info_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9), Inches(1.3),
        Inches(3.8), Inches(5.5)
    )
    info_bg.fill.solid()
    info_bg.fill.fore_color.rgb = RGBColor(30, 30, 50)
    info_bg.line.color.rgb = COLORS['secondary']
    info_bg.line.width = Pt(1)
    
    # Info title
    info_title = slide.shapes.add_textbox(Inches(9.2), Inches(1.6), Inches(3.4), Inches(0.6))
    tf = info_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 คำอธิบาย"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    
    # Info content
    info_content = slide.shapes.add_textbox(Inches(9.2), Inches(2.3), Inches(3.4), Inches(4))
    tf = info_content.text_frame
    tf.word_wrap = True
    
    info_items = [
        "setup() → ทำครั้งเดียวตอนเริ่ม",
        "loop() → ทำซ้ำไปเรื่อยๆ",
        "Serial.begin(115200) → ตั้งความเร็ว",
        "Serial.println() → พิมพ์ข้อความ",
        "delay(1000) → รอ 1 วินาที",
    ]
    
    for i, text in enumerate(info_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_light']
        p.space_after = Pt(10)

def main():
    # Create presentation (16:9 landscape)
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # 16:9
    prs.slide_height = Inches(7.5)
    
    # Create 3 slides
    create_youtube_slide_1(prs)
    create_youtube_slide_2(prs)
    create_youtube_slide_3(prs)
    
    # Save
    output_path = "/home/maxtic/esp32-stem-book/slides_yt.pptx"
    prs.save(output_path)
    print(f"✅ Created: {output_path}")
    
    # Also show file size
    size = os.path.getsize(output_path)
    print(f"📊 Size: {size/1024:.1f} KB")

if __name__ == "__main__":
    main()
